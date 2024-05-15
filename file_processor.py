from typing import Optional
import requests
import base64
import PyPDF2
import docx
from pptx import Presentation
from openpyxl import load_workbook
import csv
import os
from openai import OpenAI, AzureOpenAI

import config


class FileProcessor:

    def __init__(self, api_name='azureopenai'):
        self.api_name = api_name
        self.default_prompt = "What is the main idea of this file?"
        self.last_file_type = None
        if self.api_name == 'openai':
            self.client = OpenAI(api_key=config.openai_api_key)
            self.GPT_MODEL = config.openai_GPT_MODEL_3
            self.whisper_model = config.whisper_model_openai
        else:  # Default to 'azureopenai'
            self.api_endpoint = config.azure_api_endpoint
            self.api_version = config.azure_api_version
            self.api_key = config.azure_api_key
            self.client = AzureOpenAI(
                api_key=self.api_key,
                api_version=self.api_version,
                azure_endpoint=self.api_endpoint
            )
            self.tts_api_version = config.tts_api_version
            self.tts_model1 = config.tts_model1
            self.GPT_MODEL = config.azure_GPT_MODEL_3
            self.whisper_model = config.azure_whisper_model
            self.voice_name = config.voice_name
            self.tts_model = config.tts_model2

    def is_TunSurf_related(self, prompt):
        messages = [
        {"role": "system", "content": "Classify if the following prompt is relevant or irrelevant. Guidelines for you as a Trading Assistant:Relevance: Focus exclusively on queries related to trading and financial markets(including stock tickers). If a question falls outside this scope, politely inform the user that the question is beyond the service's focus.Accuracy: Ensure that the information provided is accurate and up-to-date. Use reliable financial data and current market analysis to inform your responses.Clarity: Deliver answers in a clear, concise, and understandable manner. Avoid jargon unless the user demonstrates familiarity with financial terms.Promptness: Aim to provide responses quickly to facilitate timely decision-making for users.Confidentiality: Do not ask for or handle personal investment details or sensitive financial information.Compliance: Adhere to legal and ethical standards applicable to financial advice and information dissemination.Again, focus solely on topics related to trading and financial markets. Politely notify the user if a question is outside this specific area of expertise."},
        {"role": "user", "content": prompt},
        ]
        response = self.client.chat.completions.create(model= self.GPT_MODEL, temperature=0.2, messages=messages)
        return response.choices[0].message.content

    def text_to_speech(self, text, output_file='output.mp3'):
      if self.api_name == 'azureopenai':
        url = f"{self.api_endpoint}/openai/deployments/{self.tts_model}/audio/speech?api-version={self.tts_api_version}"
        headers = {
            "api-key": self.api_key,
            "Content-Type": "application/json"
        }
        data = {
            "model": self.tts_model,
            "input": text,
            "voice": self.voice_name
        }
        response = requests.post(url, headers=headers, json=data)
        if response.status_code == 200:
            with open(output_file, 'wb') as f:
                f.write(response.content)
            return output_file
        else:
            print(f"Failed to generate speech: {response.status_code} - {response.text}")
            return None
      else:
        response1 = self.client.audio.speech.create(
            model= self.tts_model,
            voice= self.voice_name,
            input= text
            )
        response = response1.stream_to_file("output_file.mp3")
        if response.status_code == 200:
            with open(output_file, 'wb') as f:
                f.write(response.content)
            return output_file
        else:
            print(f"Failed to generate speech: {response.status_code} - {response.text}")
            return None

    def image_process(self, file_path: str) -> str:
        try:
            print("Image file detected. VisionGPT processing...")
            # OpenAI API Key
            api_key = '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'
            # Perform VisionGPT processing here
            def encode_image(file_path):
                with open(file_path, "rb") as image_file:
                    return base64.b64encode(image_file.read()).decode('utf-8')

            image_path = file_path
            base64_image = encode_image(image_path)

            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}"
            }

            payload = {
                "model": "gpt-4-vision-preview",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "What’s in this image?"
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{base64_image}"
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 300
            }

            response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
            response.raise_for_status()  # Raise an exception for any HTTP error status
            data_image = response.json()
            file_contents = data_image['choices'][0]['message']['content']
            print("Image processing complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the image: {e}")
            return ""

    def pdf_process(self, file_path: str) -> str:
        try:
            print("PDF file detected. Extracting text...")
            with open(file_path, 'rb') as file:
                pdf = PyPDF2.PdfReader(file)
                file_contents = ''
                for page in range(len(pdf.pages)):
                    file_contents += pdf.pages[page].extract_text()
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the PDF: {e}")
            return ""

    def text_process(self, file_path: str) -> str:
        try:
            print("Text file detected. Extracting text...")
            with open(file_path, 'r') as file:
                file_contents = file.read()
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the text file: {e}")
            return ""

    def word_process(self, file_path: str) -> str:
        try:
            print("Word document detected. Extracting text...")
            doc = docx.Document(file_path)
            file_contents = '\n'.join([para.text for para in doc.paragraphs])
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the Word document: {e}")
            return ""

    def powerpoint_process(self, file_path: str) -> str:
        try:
            print("Powerpoint document detected. Extracting text...")
            # Load the PowerPoint presentation
            prs = Presentation(file_path)
            # Initialize an empty list to store the extracted lines of text
            extracted_lines = []
            # Iterate through each slide in the presentation
            for slide in prs.slides:
                # Iterate through each shape in the slide
                for shape in slide.shapes:
                    # Check if the shape has text
                    if hasattr(shape, "text"):
                        # Strip leading and trailing spaces from the text
                        text = shape.text.strip()
                        # Append non-empty lines to the list
                        if text:
                            extracted_lines.append(text)

            # Join the extracted lines into a single string with newlines
            file_contents = '\n'.join(extracted_lines)
            print("Extraction complete.")
            # Return the extracted text
            return file_contents

        except Exception as e:
            # Print an error message if an exception occurs
            print(f"An error occurred while processing the PowerPoint file: {e}")
            return ""

    def speech_process(self, file_path: str) -> str:
        try:
            print("Speech file detected. Processing speech...")
            '''
            client = AzureOpenAI(
            api_key="80ddd1ad72504f2fa226755d49491a61",
            api_version="2023-10-01-preview",
            azure_endpoint = "https://tensurfbrain1.openai.azure.com/"
            )
            '''
            #deployment_id = "whisper_001"
            speech= open(file_path, "rb")
            transcription = self.client.audio.transcriptions.create(
            file=speech,
            model=self.whisper_model
            )
            speech_contents = transcription.text
            print("Speech processing complete.")
            return speech_contents
        except Exception as e:
            print(f"An error occurred while processing the speech file: {e}")
            return ""

    def excel_process(self, file_path: str) -> str:
        try:
            print("Excel file detected. Extracting text...")
            wb = load_workbook(file_path)
            file_contents = '\n'.join([str(cell.value) for sheet in wb.sheetnames for row in wb[sheet].iter_rows() for cell in row if cell.value is not None])
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the Excel file: {e}")
            return ""

    def csv_process(self, file_path: str) -> str:
        try:
            print("CSV file detected. Extracting text...")
            with open(file_path, 'r', newline='', encoding='utf-8') as file:
                csv_reader = csv.reader(file)
                file_contents = '\n'.join(','.join(row) for row in csv_reader)
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the CSV file: {e}")
            return ""

    def process_file(self, file_path: str) -> str:
        try:
            _, file_extension = os.path.splitext(file_path)

            self.last_file_type = file_extension.lower()  # Track the type of the last processed file

            if file_extension.lower() in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                return self.image_process(file_path)

            elif file_extension.lower() == '.pdf':
                return self.pdf_process(file_path)

            elif file_extension.lower() == '.txt':
                return self.text_process(file_path)

            elif file_extension.lower() == '.docx':
                return self.word_process(file_path)

            elif file_extension.lower() in ['.pptx', '.ppt']:
                return self.powerpoint_process(file_path)

            elif file_extension.lower() == '.mp3':
                return self.speech_process(file_path)

            elif file_extension.lower() == '.csv':
                return self.csv_process(file_path)

            elif file_extension.lower() == '.xlsx':
                return self.excel_process(file_path)

            else:
                raise ValueError("Unsupported file format. Please provide an image, PDF, PPT, TXT, DOCX, MP3, XLSX, or other supported file format.")
        except Exception as e:
            print(f"An error occurred while processing the file: {e}")
            return ""

    def get_content(self, file_path: Optional[str]) -> str:
        content = ""
        try:
            if file_path:
                file_content = self.process_file(file_path)
                content += file_content + "\n"
        except Exception as e:
            print(f"An error occurred while getting content: {e}")
        return content.strip()
