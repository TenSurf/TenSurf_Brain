import requests
import base64
import PyPDF2
import docx
from pptx import Presentation
from openpyxl import load_workbook
import csv
import os
from io import BufferedReader, StringIO

class FileProcessor:

    def __init__(self, connector_surf):
        self.default_prompt = "What is the main idea of this file?"
        self.file_type = None
        self.connector_surf = connector_surf
        
    def is_TunSurf_related(self, prompt):
        messages = [
            {
                "role": "system",
                "content": "Classify if the following prompt is relevant or irrelevant. Guidelines for you as a Trading Assistant:Relevance: Focus exclusively on queries related to trading and financial markets(including stock tickers). If a question falls outside this scope, politely inform the user that the question is beyond the service's focus.Accuracy: Ensure that the information provided is accurate and up-to-date. Use reliable financial data and current market analysis to inform your responses.Clarity: Deliver answers in a clear, concise, and understandable manner. Avoid jargon unless the user demonstrates familiarity with financial terms.Promptness: Aim to provide responses quickly to facilitate timely decision-making for users.Confidentiality: Do not ask for or handle personal investment details or sensitive financial information.Compliance: Adhere to legal and ethical standards applicable to financial advice and information dissemination.Again, focus solely on topics related to trading and financial markets. Politely notify the user if a question is outside this specific area of expertise.",
            },
            {"role": "user", "content": prompt},
        ]
        response = self.connector_surf.client.chat.completions.create(
            model=self.connector_surf.GPT_MODEL, temperature=0.2, messages=messages
        )
        return response.choices[0].message.content

    def text_to_speech(self, text):
        url = f"{self.connector_surf.api_endpoint}/openai/deployments/{self.connector_surf.tts_model}/audio/speech?api-version={self.connector_surf.tts_api_version}"
        headers = {"api-key": self.connector_surf.api_key, "Content-Type": "application/json"}
        data = {"model": self.connector_surf.tts_model, "input": text, "voice": self.connector_surf.voice_name}
        response = requests.post(url, headers=headers, json=data)
        if response.status_code == 200:
            return response.content
        else:
            print(
                f"Failed to generate speech: {response.status_code} - {response.text}"
            )
            return None

    def image_process(self, file) -> str:
        try:
            print("Image file detected. VisionGPT processing...")
            # OpenAI API Key
            api_key = "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"

            base64_image = base64.b64encode(file.read()).decode("utf-8")

            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}",
            }

            payload = {
                "model": "gpt-4-vision-preview",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "What’s in this image?"},
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{base64_image}"
                                },
                            },
                        ],
                    }
                ],
                "max_tokens": 300,
            }

            response = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers=headers,
                json=payload,
            )
            response.raise_for_status()  # Raise an exception for any HTTP error status
            data_image = response.json()
            file_contents = data_image["choices"][0]["message"]["content"]
            print("Image processing complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the image: {e}")
            return ""

    def pdf_process(self, file) -> str:
        try:
            print("PDF file detected. Extracting text...")
            pdf = PyPDF2.PdfReader(file)
            file_contents = ""
            for page in range(len(pdf.pages)):
                file_contents += pdf.pages[page].extract_text()
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the PDF: {e}")
            return ""

    def text_process(self, file) -> str:
        try:
            print("Text file detected. Extracting text...")
            file_contents = file.read().decode("utf-8")
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the text file: {e}")
            return ""

    def word_process(self, file) -> str:
        try:
            print("Word document detected. Extracting text...")
            doc = docx.Document(file)
            file_contents = "\n".join([para.text for para in doc.paragraphs])
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the Word document: {e}")
            return ""

    def powerpoint_process(self, file) -> str:
        try:
            print("Powerpoint document detected. Extracting text...")
            # Load the PowerPoint presentation
            prs = Presentation(file)
            # Initialize an empty list to store the extracted lines of text
            extracted_lines = []
            # Iterate through each slide in the presentation
            for slide in prs.slides:
                # Iterate through each shape in the slide
                for shape in slide.shapes:
                    # Check if the shape has text
                    if hasattr(shape, "text"):
                        # Strip leading and trailing spaces from the text
                        text = shape.text.strip()
                        # Append non-empty lines to the list
                        if text:
                            extracted_lines.append(text)

            # Join the extracted lines into a single string with newlines
            file_contents = "\n".join(extracted_lines)
            print("Extraction complete.")
            # Return the extracted text
            return file_contents

        except Exception as e:
            # Print an error message if an exception occurs
            print(f"An error occurred while processing the PowerPoint file: {e}")
            return ""

    def speech_process(self, file) -> str:
        try:
            print("Speech file detected. Processing speech...")
            """
            client = AzureOpenAI(
            api_key="80ddd1ad72504f2fa226755d49491a61",
            api_version="2023-10-01-preview",
            azure_endpoint = "https://tensurfbrain1.openai.azure.com/"
            )
            """
            # deployment_id = "whisper_001"
            transcription = self.connector_surf.client.audio.transcriptions.create(
                file=BufferedReader(file), model=self.connector_surf.whisper_model
            )
            speech_contents = transcription.text
            print("Speech processing complete.")
            return speech_contents
        except Exception as e:
            print(f"An error occurred while processing the speech file: {e}")
            return ""

    def excel_process(self, file) -> str:
        try:
            print("Excel file detected. Extracting text...")
            wb = load_workbook(file)
            file_contents = "\n".join(
                [
                    str(cell.value)
                    for sheet in wb.sheetnames
                    for row in wb[sheet].iter_rows()
                    for cell in row
                    if cell.value is not None
                ]
            )
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the Excel file: {e}")
            return ""

    def csv_process(self, file) -> str:
        try:
            print("CSV file detected. Extracting text...")
            csv_reader = csv.reader(StringIO(file.read().decode("utf-8")))
            file_contents = "\n".join(",".join(row) for row in csv_reader)
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the CSV file: {e}")
            return ""

    def process_file(self, file) -> str:
        try:
            _, file_extension = os.path.splitext(file.name)

            self.file_type = (
                file_extension.lower()
            )  # Track the type of the last processed file

            if self.file_type in [
                ".jpg",
                ".jpeg",
                ".png",
                ".gif",
                ".bmp",
                ".tiff",
            ]:
                return self.image_process(file)

            elif self.file_type == ".pdf":
                return self.pdf_process(file)

            elif self.file_type == ".txt":
                return self.text_process(file)

            elif self.file_type == ".docx":
                return self.word_process(file)

            elif self.file_type in [".pptx", ".ppt"]:
                return self.powerpoint_process(file)

            elif self.file_type in [".mp3", ".wav"]:
                return self.speech_process(file)

            elif self.file_type == ".csv":
                return self.csv_process(file)

            elif self.file_type == ".xlsx":
                return self.excel_process(file)

            else:
                raise ValueError(
                    "Unsupported file format. Please provide an image, PDF, PPT, TXT, DOCX, MP3, XLSX, or other supported file format."
                )
        except Exception as e:
            print(f"An error occurred while processing the file: {e}")
            return ""

    def get_file_content(self, file) -> str:
        try:
            if file:
                return self.process_file(file)
            return ""
        except Exception as e:
            print(f"An error occurred while getting content: {e}")
            return ""
