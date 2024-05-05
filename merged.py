# -*- coding: utf-8 -*-
"""Final_FileProcessing_GPT.ipynb

Automatically generated by Colab.

Original file is located at
    https://colab.research.google.com/drive/1o_VWgXzFbICtCyLz76eeH1LYy6A6tZWF
"""

from openai import OpenAI, AzureOpenAI
import PyPDF2
import os
from typing import Optional
import requests
from typing import Optional
import base64
from openpyxl import load_workbook
import requests
import json
import csv
from pptx import Presentation
import docx
from dateutil.relativedelta import relativedelta
from openai import AzureOpenAI

from gpt.functions_json import functions
from gpt.functions_python import *
from gpt.utils import date_validation, monthdelta
from io import BufferedReader, StringIO
from datetime import timezone, datetime, timedelta
from gpt.input_filter import *


class FileProcessor:
    def __init__(self, api_name='azureopenai'):
        if api_name == 'openai':
            self.api_key = "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"
            self.client = OpenAI(api_key=self.api_key)
            self.GPT_MODEL = "gpt-3.5-turbo-1106"
            self.whisper_model = "whisper-1"
        else:  # Default to 'azureopenai'
            self.api_endpoint = 'https://tensurfbrain1.openai.azure.com/'
            self.api_version = '2023-10-01-preview'
            self.api_key = '80ddd1ad72504f2fa226755d49491a61'
            self.client = AzureOpenAI(
                api_key=self.api_key,
                api_version=self.api_version,
                azure_endpoint=self.api_endpoint
            )
            self.tts_api_version = '2024-02-15-preview'
            self.tts_model = 'TTS_1'
            self.GPT_MODEL = "gpt_35_16k"
            self.whisper_model = "whisper_001"
            self.voice_name = 'alloy'
            self.tts_model = 'tts_1'
            
    def is_TunSurf_related (self, prompt):
        messages = [
            {"role": "system", "content": "Classify if the following prompt is relevant or irrelevant. Guidelines for you as a Trading Assistant:Relevance: Focus exclusively on queries related to trading and financial markets. If a question falls outside this scope, politely inform the user that the question is beyond the service's focus.Accuracy: Ensure that the information provided is accurate and up-to-date. Use reliable financial data and current market analysis to inform your responses.Clarity: Deliver answers in a clear, concise, and understandable manner. Avoid jargon unless the user demonstrates familiarity with financial terms.Promptness: Aim to provide responses quickly to facilitate timely decision-making for users.Confidentiality: Do not ask for or handle personal investment details or sensitive financial information.Compliance: Adhere to legal and ethical standards applicable to financial advice and information dissemination.Again, focus solely on topics related to trading and financial markets. Politely notify the user if a question is outside this specific area of expertise."},
            {"role": "user", "content": prompt},
        ]
        response = self.client.chat.completions.create(model= self.GPT_MODEL ,messages=messages)
        #print(response)
        return response.choices[0].message.content
    
    def text_to_speech(self, text, output_file='output.mp3'):
      if self.api_name == 'azureopenai':
        url = f"{self.api_endpoint}/openai/deployments/{self.tts_model}/audio/speech?api-version={self.tts_api_version}"
        headers = {
            "api-key": self.api_key,
            "Content-Type": "application/json"
        }
        data = {
            "model": self.tts_model,
            "input": text,
            "voice": self.voice_name
        }
        response = requests.post(url, headers=headers, json=data)
        if response.status_code == 200:
            with open(output_file, 'wb') as f:
                f.write(response.content)
            return output_file
        else:
            print(f"Failed to generate speech: {response.status_code} - {response.text}")
            return None
      else:
        response1 = self.client.audio.speech.create(
            model= self.tts_model,
            voice= self.voice_name,
            input= text
            )
        response = response1.stream_to_file("output_file.mp3")
        if response.status_code == 200:
            with open(output_file, 'wb') as f:
                f.write(response.content)
            return output_file
        else:
            print(f"Failed to generate speech: {response.status_code} - {response.text}")
            return None

    def chat_with_ai(self, messages: list, content: str, front_json: dict, file_exist):
        if not content:
            return ""
        
        relevance_check = self.is_TunSurf_related(content)
        if "irrelevant" in relevance_check.lower():
            return "I'm here to help with trading and financial market queries. If you think your ask relates to trading and isn't addressed, please report a bug using the bottom right panel."
        
        def get_response(messages, functions, model, function_call, temperature=0.2):
            response = self.client.chat.completions.create(
                model=model, messages=messages, functions=functions, function_call=function_call, temperature=temperature)
            return response
        
        def get_result(messages, chat_response):
            assistant_message = chat_response.choices[0].message
            messages.append(assistant_message)
            
            if chat_response.choices[0].message.function_call == None:
                results = f"{chat_response.choices[0].message.content}"
            
            else:
                results = ""
                if assistant_message.content != None:
                    results += assistant_message.content + "\n"
                function_name = chat_response.choices[0].message.function_call.name
                function_arguments = json.loads(chat_response.choices[0].message.function_call.arguments)
                FC = FunctionCalls()
                print(f"\n{chat_response.choices[0].message}\n")
                now = datetime.now()

                # Filtering Inputs
                function_arguments = input_filter(function_name, function_arguments, front_json)
                if function_name == "detect_trend":
                    # correct_dates = True
                    # # validating dates
                    # if "start_datetime" in function_arguments or "end_datetime" in function_arguments:
                    #     correct_dates = False
                    #     # checking the format
                    #     if not (date_validation(function_arguments["start_datetime"]) or date_validation(function_arguments["end_datetime"])):
                    #         results += "Please enter dates in the following foramat: %d/%m/%Y %H:%M:%S or specify a period of time whether for the past seconds or minutes or hours or days or weeks or years."
                    #     # if start_datetime or end_datetime were in the future
                    #     elif (now < datetime.strptime(function_arguments["start_datetime"], '%d/%m/%Y %H:%M:%S')) or (now < datetime.strptime(function_arguments["end_datetime"], '%d/%m/%Y %H:%M:%S')):
                    #         results += "Dates should not be in the future!"
                    #     # if end is before start
                    #     elif datetime.strptime(function_arguments["end_datetime"], '%d/%m/%Y %H:%M:%S') < datetime.strptime(function_arguments["start_datetime"], '%d/%m/%Y %H:%M:%S'):
                    #         results += "End date time should be after start date time!"
                    #     # formates are correct
                    #     elif "lookback" in function_arguments and "start_datetime" in function_arguments:
                    #         results += "Both lookback and datetimes could not be valued"
                    #     # dates are valid
                    #     else:
                    #         correct_dates = True
                    
                    # # if loookback and end_datetime were specified
                    # if ("lookback" in function_arguments and "end_datetime" in function_arguments) and correct_dates:
                    #     end_datetime = datetime.strptime(function_arguments["end_datetime"], '%d/%m/%Y %H:%M:%S')
                    #     k = int(function_arguments["lookback"].split(" ")[0])
                    #     function_arguments["start_datetime"] = f"{end_datetime - timedelta(days=k)}"
                    
                    # # handling the default values when none of the parameters were specified
                    # elif ("lookback" not in function_arguments) and correct_dates:
                    #     if "symbol" not in function_arguments:
                    #         function_arguments["symbol"] = "NQ"
                    #     if "start_datetime" not in function_arguments:
                    #         function_arguments["start_datetime"] = f"{now - timedelta(days=10)}"
                    #     if "end_datetime" not in function_arguments:
                    #         function_arguments["end_datetime"] = f"{now}"
                    
                    # # if just lookback was specified
                    # elif correct_dates:
                    #     function_arguments["end_datetime"] = f"{now}"
                    #     k = int(function_arguments["lookback"].split(" ")[0])
                    #     if function_arguments["lookback"].split(" ")[-1] == "seconds" or function_arguments["lookback"].split(" ")[-1] == "second":
                    #         function_arguments["start_datetime"] = f"{now - timedelta(seconds=k)}"
                    #     elif function_arguments["lookback"].split(" ")[-1] == "minutes" or function_arguments["lookback"].split(" ")[-1] == "minute":
                    #         function_arguments["start_datetime"] = f"{now - timedelta(minutes=k)}"
                    #     elif function_arguments["lookback"].split(" ")[-1] == "hours" or function_arguments["lookback"].split(" ")[-1] == "hour":
                    #         function_arguments["start_datetime"] = f"{now - timedelta(hours=k)}"
                    #     elif function_arguments["lookback"].split(" ")[-1] == "days" or function_arguments["lookback"].split(" ")[-1] == "day":
                    #         function_arguments["start_datetime"] = f"{now - timedelta(days=k)}"
                    #     elif function_arguments["lookback"].split(" ")[-1] == "weeks" or function_arguments["lookback"].split(" ")[-1] == "week":
                    #         function_arguments["start_datetime"] = f"{now - timedelta(weeks=k)}"
                    #     elif function_arguments["lookback"].split(" ")[-1] == "months" or function_arguments["lookback"].split(" ")[-1] == "month":
                    #         function_arguments["start_datetime"] = f"{monthdelta(now, -k)}"
                    #     elif function_arguments["lookback"].split(" ")[-1] == "years" or function_arguments["lookback"].split(" ")[-1] == "year":
                    #         function_arguments["start_datetime"] = f"{now - relativedelta(years=k)}"
                    #     else:
                    #         raise ValueError("wrong value of time")
                    
                    function_arguments, results, correct_dates = function_arguments

                    # results will be generated only when dates are in the correct format
                    if correct_dates:
                        trend = FC.detect_trend(parameters=function_arguments)
                        messages.append({"role": "system", "content": f"The result of the function calling with function {function_name} has become {trend}. At any situations, never return the number which is the output of the detect_trend function. Instead, use its correcsponding explanation which is in the detect_trend function's description. Make sure to mention the start_datetime and end_datetime. If the user provide neither specified both start_datetime and end_datetime nor lookback parameters, politely tell them that they should and introduce these parameters to them so that they can use them. Do not mention the name of the parameters of the functions directly in the final answer. Instead, briefly explain them and use other meaningfuly related synonyms. Now generate a proper response."})
                        chat_response = get_response(
                            messages, functions, self.GPT_MODEL, "auto"
                        )
                        assistant_message = chat_response.choices[0].message
                        messages.append(assistant_message)
                        results += chat_response.choices[0].message.content

                elif function_name == "calculate_sr":
                    # correcting function_arguments
                    # if "symbol" not in function_arguments:
                    #     function_arguments["symbol"] = "ES"
                    # if "timeframe" not in function_arguments:
                    #     function_arguments["timeframe"] = "1h"
                    # if "lookback_days" not in function_arguments:
                    #     function_arguments["lookback_days"] = "10 days"

                    sr_value, sr_start_date, sr_end_date, sr_importance = FC.calculate_sr(parameters=function_arguments)
                    timezone_number = int(front_json["timezone"])
                    for date in sr_start_date:
                        date -= timedelta(minutes=timezone_number)
                    for date in sr_end_date:
                        date -= timedelta(minutes=timezone_number)
                    messages.append({"role": "system", "content": f"The result of the function calling with function {function_name} has become {sr_value} for levels_prices, {sr_start_date} for levels_start_timestamps, {sr_end_date} for levels_end_timestamps and {sr_importance} for levels_scores. Do not mention the name of the parameters of the functions directly in the final answer. Instead, briefly explain them and use other meaningfuly related synonyms. If the user didn't specified lookback_days or timeframe parameters, introduce these parameters to them so that they can use these parameters. Now generate a proper response"})
                    chat_response = get_response(
                        messages, functions, self.GPT_MODEL, "auto"
                    )
                    results = {
                        'content': results + chat_response.choices[0].message.content,
                        'levels': { 'value': sr_value, 'start': sr_start_date, 'end': sr_end_date, 'importance': sr_importance }
                    }
                    

                elif function_name == "calculate_sl":
                    stoploss = FC.calculate_sl(function_arguments)
                    messages.append({"role": "system", "content": f"The result of the function calling with function {function_name} has become {stoploss}. Do not mention the name of the parameters of the functions directly in the final answer. Instead, briefly explain them and use other meaningfuly related synonyms. The unit of every number in the answer should be mentioned. Now generate a proper response"})
                    chat_response = get_response(
                        messages, functions, self.GPT_MODEL, "auto"
                    )
                    results += chat_response.choices[0].message.content

                elif function_name == "calculate_tp":
                    takeprofit = FC.calculate_tp(function_arguments)
                    messages.append({"role": "system", "content": f"The result of the function calling with function {function_name} has become {takeprofit}. Do not mention the name of the parameters of the functions directly in the final answer. Instead, briefly explain them and use other meaningfuly related synonyms. Now generate a proper response"})
                    chat_response = get_response(
                        messages, functions, self.GPT_MODEL, "auto"
                    )
                    results += chat_response.choices[0].message.content
                    

                else:
                    results += f"{chat_response.choices[0].message.content}"
                
            return results
        
        messages.append({"role": "system", "content": "Don't make assumptions about what values to plug into functions. Ask for clarification if a user request is ambiguous. Make sure to restrict your function calls to the following provided list: detect_trend, calculate_sr, calculate_sl, calculate_tp. Do not assume the availability of any functions beyond those explicitly listed. Ensure that your implementations and queries adhere strictly to the functions specified."})
        messages.append({"role": "user", "content": content})
        response = get_response(messages, functions, self.GPT_MODEL, "auto")
        res = get_result(messages, response)
        if file_exist ==1:
            if self.last_file_type == '.mp3':  # Check if the last file processed was MP3
                self.text_to_speech(res, 'response.mp3')
                return res, 'response.mp3'
        return res
    
        messages.append({"role": "system", "content": "Don't make assumptions about what values to plug into functions. Ask for clarification if a user request is ambiguous."})
        messages.append({"role": "user", "content": content})
        response = get_response(messages, functions, self.GPT_MODEL, "auto")
        res = get_result(messages, response)
        if file_exist ==1:
            if self.last_file_type == '.mp3':  # Check if the last file processed was MP3
                self.text_to_speech(res, 'response.mp3')
                return res, 'response.mp3'
        return res

    def image_process(self, file) -> str:
        try:
            print("Image file detected. VisionGPT processing...")
            # OpenAI API Key
            api_key = os.environ.get("OPENAI_API_KEY")
            # Perform VisionGPT processing here
            def encode_image(file):
                return base64.b64encode(file.read()).decode('utf-8')

            base64_image = encode_image(file)

            headers = {
                "Content-Type": "application/json",
                "Authorization": f"Bearer {self.openai_api_key}"
            }

            payload = {
                "model": "gpt-4-vision-preview",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "What’s in this image?"
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{base64_image}"
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 300
            }

            response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
            response.raise_for_status()  # Raise an exception for any HTTP error status
            data_image = response.json()
            file_contents = data_image['choices'][0]['message']['content']
            print("Image processing complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the image: {e}")
            return ""

    def pdf_process(self, file) -> str:
        try:
            print("PDF file detected. Extracting text...")
            pdf = PyPDF2.PdfReader(file)
            file_contents = ''
            for page in range(len(pdf.pages)):
                file_contents += pdf.pages[page].extract_text()
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the PDF: {e}")
            return ""

    def text_process(self, file) -> str:
        try:
            print("Text file detected. Extracting text...")
            file_contents = file.read().decode('utf-8')
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the text file: {e}")
            return ""

    def word_process(self, file) -> str:
        try:
            print("Word document detected. Extracting text...")
            doc = docx.Document(file)
            file_contents = '\n'.join([para.text for para in doc.paragraphs])
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the Word document: {e}")
            return ""

    def powerpoint_process(self, file) -> str:
        try:
            print("Powerpoint document detected. Extracting text...")
            # Load the PowerPoint presentation
            prs = Presentation(file)
            # Initialize an empty list to store the extracted lines of text
            extracted_lines = []
            # Iterate through each slide in the presentation
            for slide in prs.slides:
                # Iterate through each shape in the slide
                for shape in slide.shapes:
                    # Check if the shape has text
                    if hasattr(shape, "text"):
                        # Strip leading and trailing spaces from the text
                        text = shape.text.strip()
                        # Append non-empty lines to the list
                        if text:
                            extracted_lines.append(text)

            # Join the extracted lines into a single string with newlines
            file_contents = '\n'.join(extracted_lines)
            print("Extraction complete.")
            # Return the extracted text
            return file_contents

        except Exception as e:
            # Print an error message if an exception occurs
            print(f"An error occurred while processing the PowerPoint file: {e}")
            return ""

    def speech_process(self, speech) -> str:
        try:
            print("Speech file detected. Processing speech...")
            transcription = self.client.audio.transcriptions.create(
                model=self.whisper_model,
                file= BufferedReader(speech))
            speech_contents = transcription.text
            print("Speech processing complete.")
            return speech_contents
        except Exception as e:
            print(f"An error occurred while processing the speech file: {e}")
            return ""

    def excel_process(self, file) -> str:
        try:
            print("Excel file detected. Extracting text...")
            wb = load_workbook(file)
            file_contents = '\n'.join([str(cell.value) for sheet in wb.sheetnames for row in wb[sheet].iter_rows() for cell in row if cell.value is not None])
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the Excel file: {e}")
            return ""

    def csv_process(self, file) -> str:
        try:
            print("CSV file detected. Extracting text...")
            csv_reader = csv.reader(StringIO(file.read().decode('utf-8')))
            file_contents = '\n'.join(','.join(row) for row in csv_reader)
            print("Extraction complete.")
            return file_contents
        except Exception as e:
            print(f"An error occurred while processing the CSV file: {e}")
            return ""

    def process_file(self, file) -> str:
        try:
            _, file_extension = os.path.splitext(file.name)
            
            self.last_file_type = file_extension.lower()  # Track the type of the last processed file

            if file_extension.lower() in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                return self.image_process(file)

            elif file_extension.lower() == '.pdf':
                return self.pdf_process(file)

            elif file_extension.lower() == '.txt':
                return self.text_process(file)

            elif file_extension.lower() == '.docx':
                return self.word_process(file)

            elif file_extension.lower() in ['.pptx', '.ppt']:
                return self.powerpoint_process(file)

            elif file_extension.lower() in ['.mp3', '.wav']:
                return self.speech_process(file)

            elif file_extension.lower() == '.csv':
                return self.csv_process(file)

            elif file_extension.lower() == '.xlsx':
                return self.excel_process(file)

            else:
                raise ValueError("Unsupported file format. Please provide an image, PDF, PPT, TXT, DOCX, MP3, XLSX, or other supported file format.")
        except Exception as e:
            print(f"An error occurred while processing the file: {e}")
            return ""

    def get_content(self, file) -> str:
        content = ""
        try:
            if file:
                file_content = self.process_file(file)
                content += file_content + "\n"
        except Exception as e:
            print(f"An error occurred while getting content: {e}")
        return content.strip()

    def get_user_input(self, file, prompt: Optional[str], messages: Optional[list], front_json: Optional[dict]) -> str:
        content = ""
        file_exist = 0
        if file:
            file_exist = 1
            file_content = self.get_content(file)
            if file_content:
                content += file_content + "\n"
        if prompt:
            content += prompt + "\n"
        else:
            content += 'what is the main idea?' + '\n'
        return self.chat_with_ai(messages=messages,content=content.strip(), front_json=front_json, file_exist=file_exist)
